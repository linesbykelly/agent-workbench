from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Optional, Dict
from pydantic import BaseModel
from pathlib import Path
from agents import function_tool
import logging

from utils import _get_files_path


logger = logging.getLogger(__name__)


class GridCell(BaseModel):
    text: Optional[str] = None
    image: Optional[str] = None


def create_file(filename: str) -> Dict:
    try:
        file_path = _get_files_path(filename)

        prs = Presentation()
        prs.save(file_path)

        return {
            "success": True,
            "data": f"Presentation created and saved at {file_path}",
            "error": None
        }

    except (OSError, ValueError) as e:
        logger.exception(
            "Failed to create presentation: %s",
            filename
        )

        return {
            "success": False,
            "data": None,
            "error": {
                "type": type(e).__name__,
                "message": str(e)
            }
        }


@function_tool
def create_pptx_file(filename: str) -> Dict:
    """
    Create a new PowerPoint presentation in the files directory.

    Args:
        filename: Filename for the PowerPoint file.
            Example: "my_presentation.pptx"

    Returns:
        A result indicating whether the presentation was created.
    """
    return create_file(filename)


def add_slide(
    filename: str,
    title: str,
    content: str
) -> Dict:

    try:
        file_path = _get_files_path(filename)

        prs = Presentation(file_path)

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        prs.save(file_path)

        return {
            "success": True,
            "data": f"Slide added with title: '{title}'",
            "error": None
        }

    except (OSError, ValueError) as e:
        logger.exception(
            "Failed to add slide to: %s",
            filename
        )

        return {
            "success": False,
            "data": None,
            "error": {
                "type": type(e).__name__,
                "message": str(e)
            }
        }


@function_tool
def add_pptx_slide(
    filename: str,
    title: str,
    content: str
) -> Dict:
    """
    Add a new slide with a title and content to an existing
    PowerPoint file in the files directory.

    Args:
        filename: Filename of the existing PowerPoint file.
        title: Title of the slide.
        content: Body content of the slide.

    Returns:
        A result indicating whether the slide was added.
    """
    return add_slide(
        filename,
        title,
        content
    )


def add_grid_slide(
    filename: str,
    title: str,
    subtitle: str,
    grid: list[list[GridCell]],
    margin: float = 0.5,
    cell_width: float = 3.0,
    cell_height: float = 2.0,
    spacing: float = 0.25
) -> Dict:

    try:
        file_path = _get_files_path(filename)

        prs = Presentation(file_path)

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        slide.shapes.title.text = title

        # Subtitle
        left = Inches(margin)
        top = Inches(1.0)
        width = prs.slide_width - Inches(2 * margin)
        height = Inches(0.5)

        subtitle_box = slide.shapes.add_textbox(
            left,
            top,
            width,
            height
        )

        subtitle_box.text_frame.text = subtitle

        # Validate grid
        rows = len(grid)
        cols = len(grid[0]) if rows > 0 else 0

        if rows == 0 or cols == 0:
            raise ValueError("Grid cannot be empty.")

        if any(len(row) != cols for row in grid):
            raise ValueError(
                "All grid rows must have the same number of columns."
            )

        # Add grid cells
        for r in range(rows):
            for c in range(cols):

                cell = grid[r][c]

                logger.debug(
                    "Cell (%d, %d): text=%r, image=%r",
                    r,
                    c,
                    cell.text,
                    cell.image
                )

                x = Inches(
                    margin + c * (cell_width + spacing)
                )

                y = Inches(
                    1.5 + r * (cell_height + spacing)
                )

                # Image
                if cell.image:

                    image_path = _get_files_path(
                        cell.image
                    )

                    if not image_path.exists():
                        raise FileNotFoundError(
                            f"Image file not found: {cell.image}"
                        )

                    slide.shapes.add_picture(
                        str(image_path),
                        x,
                        y,
                        width=Inches(cell_width),
                        height=Inches(cell_height)
                    )

                # Text
                if cell.text:

                    text_box = slide.shapes.add_textbox(
                        x,
                        y,
                        Inches(cell_width),
                        Inches(cell_height)
                    )

                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True

                    paragraph = text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER

                    run = paragraph.add_run()
                    run.text = cell.text

                    font = run.font
                    font.size = Pt(12)
                    font.name = "Calibri"
                    font.color.rgb = RGBColor(0, 0, 0)

        prs.save(file_path)

        return {
            "success": True,
            "data": f'Grid slide added with title: "{title}"',
            "error": None
        }

    except (OSError, ValueError, FileNotFoundError) as e:
        logger.exception(
            "Failed to add grid slide to: %s",
            filename
        )

        return {
            "success": False,
            "data": None,
            "error": {
                "type": type(e).__name__,
                "message": str(e)
            }
        }


@function_tool
def add_pptx_slide_w_grid(
    filename: str,
    title: str,
    subtitle: str,
    grid: list[list[GridCell]],
    margin: float = 0.5,
    cell_width: float = 3.0,
    cell_height: float = 2.0,
    spacing: float = 0.25
) -> Dict:
    """
    Add a slide with a title, subtitle, and customizable grid
    of text and/or images.

    Args:
        filename: Filename of the existing PowerPoint file.
        title: Title displayed at the top of the slide.
        subtitle: Subtitle displayed below the title.
        grid: 2D list of GridCell objects.
        margin: Margin from slide borders in inches.
        cell_width: Width of each cell in inches.
        cell_height: Height of each cell in inches.
        spacing: Spacing between cells in inches.

    Returns:
        A result indicating whether the grid slide was added.
    """

    return add_grid_slide(
        filename,
        title,
        subtitle,
        grid,
        margin,
        cell_width,
        cell_height,
        spacing
    )


def read_file_text(filename: str) -> Dict:

    try:
        file_path = _get_files_path(filename)

        prs = Presentation(file_path)

        all_text = []

        for i, slide in enumerate(prs.slides, 1):

            slide_text = f"Slide {i}:\n"

            for shape in slide.shapes:

                if (
                    hasattr(shape, "text")
                    and shape.text.strip()
                ):
                    slide_text += shape.text + "\n"

            all_text.append(
                slide_text.strip()
            )

        return {
            "success": True,
            "data": "\n\n".join(all_text),
            "error": None
        }

    except (OSError, ValueError) as e:
        logger.exception(
            "Failed to read presentation: %s",
            filename
        )

        return {
            "success": False,
            "data": None,
            "error": {
                "type": type(e).__name__,
                "message": str(e)
            }
        }


@function_tool
def read_pptx_file_text(filename: str) -> Dict:
    """
    Read and extract all text content from a PowerPoint presentation.

    Args:
        filename: Filename of the PowerPoint file.

    Returns:
        A result containing the text organized by slide number.
    """
    return read_file_text(filename)